import os
from typing import Any, Dict

from utils.logger import get_logger

logger = get_logger()

# Supported file formats and their handlers
SUPPORTED_EXTENSIONS = {
    # Documents
    ".pdf": "_read_pdf",
    ".docx": "_read_docx",
    ".doc": "_read_docx",
    ".pptx": "_read_pptx",
    ".xlsx": "_read_excel",
    ".xls": "_read_excel",
    ".csv": "_read_csv",
    
    # Notebooks
    ".ipynb": "_read_ipynb",

    # Plain Text / Code / Config
    ".txt": "_read_text",
    ".md": "_read_text",
    ".json": "_read_text",
    ".jsonc": "_read_text",
    ".json5": "_read_text",
    ".yaml": "_read_text",
    ".yml": "_read_text",
    ".log": "_read_text",
    ".py": "_read_text",
    ".js": "_read_text",
    ".mjs": "_read_text",
    ".ts": "_read_text",
    ".tsx": "_read_text",
    ".jsx": "_read_text",
    ".html": "_read_text",
    ".css": "_read_text",
    ".scss": "_read_text",
    ".sass": "_read_text",
    ".go": "_read_text",
    ".rs": "_read_text",
    ".cpp": "_read_text",
    ".c": "_read_text",
    ".h": "_read_text",
    ".hpp": "_read_text",
    ".java": "_read_text",
    ".kt": "_read_text",
    ".kts": "_read_text",
    ".swift": "_read_text",
    ".rb": "_read_text",
    ".php": "_read_text",
    ".lua": "_read_text",
    ".dart": "_read_text",
    ".sql": "_read_text",
    ".sh": "_read_text",
    ".bash": "_read_text",
    ".bat": "_read_text",
    ".ps1": "_read_text",
    ".env": "_read_text",
    ".ini": "_read_text",
    ".toml": "_read_text",
    ".conf": "_read_text",
    ".xml": "_read_text",
    ".props": "_read_text",
    ".properties": "_read_text",
    ".gitignore": "_read_text",
    ".dockerfile": "_read_text"
}

TOOL_META = {
    "name": "file_reader",
    "aliases": ["read_file", "read_document", "read_pdf", "read_excel", "read_csv", "read_code", "read_notebook"],
    "class_name": "FileReader",
    "description": f"File Reader: Extracts text and structural data from documents, notebooks, and source code. Supported extensions: {', '.join(sorted(SUPPORTED_EXTENSIONS.keys()))}",
    "actions": [
        {
            "name": "read",
            "description": "Read and extract content from a document or code file. Supports modern formats like .ipynb, .pptx and various source code files.",
            "input_schema": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the document file."},
                    "max_pages": {"type": "integer", "description": "Max pages to read (PDF). Default: 50."}
                },
                "required": ["file_path"]
            }
        }
    ]
}


class FileReader:
    """Extracts content from various document formats using the global SUPPORTED_EXTENSIONS mapping."""

    @classmethod
    async def execute(cls, file_path: str = "", max_pages: int = 50, **kwargs) -> Dict[str, Any]:
        """Read a document file and extract its content."""
        if not file_path:
            return {"status": "error", "message": "Missing 'file_path' parameter."}

        if not os.path.exists(file_path):
            return {"status": "error", "message": f"File not found: {file_path}"}

        _, filename = os.path.split(file_path)
        if filename.startswith(".") and not os.path.splitext(filename)[1]:
            # It's a dotfile like .env, treat the whole thing as extension
            ext = filename.lower()
        else:
            ext = os.path.splitext(file_path)[1].lower()
            
        handler_name = SUPPORTED_EXTENSIONS.get(ext)

        if not handler_name:
            return {
                "status": "error",
                "message": f"Unsupported file format: {ext}. Total supported: {len(SUPPORTED_EXTENSIONS)} extensions."
            }

        try:
            handler = getattr(cls, handler_name)
            content = handler(file_path, max_pages=max_pages)
            logger.info(f"[FILE_READER] Read {ext} file: {file_path} ({len(content)} chars)")
            return {
                "status": "success",
                "data": {
                    "content": content,
                    "file_name": os.path.basename(file_path),
                    "file_type": ext,
                    "char_count": len(content)
                }
            }
        except Exception as e:
            logger.error(f"[FILE_READER] Failed to read {file_path}: {e}")
            return {"status": "error", "message": f"Failed to read file: {str(e)}"}

    @staticmethod
    def _read_pdf(file_path: str, max_pages: int = 50, **kwargs) -> str:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        pages = reader.pages[:max_pages]
        text_parts = []
        for i, page in enumerate(pages):
            text = page.extract_text() or ""
            if text.strip():
                text_parts.append(f"--- Page {i+1} ---\n{text.strip()}")
        return "\n\n".join(text_parts) if text_parts else "No text could be extracted from this PDF."

    @staticmethod
    def _read_docx(file_path: str, **kwargs) -> str:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        # Also extract tables
        table_texts = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            table_texts.append("\n".join(rows))
        
        all_content = "\n".join(paragraphs)
        if table_texts:
            all_content += "\n\n--- Tables ---\n" + "\n\n".join(table_texts)
        return all_content or "No text could be extracted from this document."

    @staticmethod
    def _read_pptx(file_path: str, **kwargs) -> str:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
        return "\n\n".join(text_parts) if text_parts else "No text could be extracted from this PowerPoint."

    @staticmethod
    def _read_excel(file_path: str, **kwargs) -> str:
        import pandas as pd
        xls = pd.ExcelFile(file_path)
        parts = []
        for sheet in xls.sheet_names[:10]:  # Limit to 10 sheets
            df = pd.read_excel(xls, sheet_name=sheet, nrows=500)
            parts.append(f"--- Sheet: {sheet} ({len(df)} rows) ---\n{df.to_markdown(index=False)}")
        return "\n\n".join(parts)

    @staticmethod
    def _read_csv(file_path: str, **kwargs) -> str:
        import pandas as pd
        df = pd.read_csv(file_path, nrows=1000)
        return f"CSV ({len(df)} rows, {len(df.columns)} columns):\n{df.to_markdown(index=False)}"

    @staticmethod
    def _read_ipynb(file_path: str, **kwargs) -> str:
        import nbformat
        with open(file_path, "r", encoding="utf-8") as f:
            nb = nbformat.read(f, as_version=4)
        
        parts = []
        for i, cell in enumerate(nb.cells):
            cell_type = cell.get("cell_type", "unknown")
            source = cell.get("source", "").strip()
            if not source:
                continue
                
            if cell_type == "markdown":
                parts.append(f"### [Markdown Cell {i+1}]\n{source}")
            elif cell_type == "code":
                parts.append(f"### [Code Cell {i+1}]\n```python\n{source}\n```")
                # Optionally include output summary
                outputs = cell.get("outputs", [])
                for out in outputs:
                    if out.get("output_type") == "stream" and "text" in out:
                        parts.append(f"**Output:**\n{out['text'].strip()}")
                    elif out.get("output_type") == "execute_result" and "data" in out:
                        if "text/plain" in out["data"]:
                            parts.append(f"**Result:** {out['data']['text/plain']}")
                            
        return "\n\n".join(parts) if parts else "No content found in this notebook."

    @staticmethod
    def _read_text(file_path: str, **kwargs) -> str:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read(500_000)  # 500KB limit
